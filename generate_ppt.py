"""Generate a PPT explaining finviz_screener_new_high.py"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT    = RGBColor(0x16, 0x21, 0x3E)   # panel bg
HIGHLIGHT = RGBColor(0x0F, 0x3C, 0x78)   # blue
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # positive
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)   # mid tier
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xCC, 0xCC, 0xCC)
RED       = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_bullet_box(slide, title, bullets, l, t, w, h,
                   title_color=GREEN, title_size=16, bullet_size=13,
                   bg_color=ACCENT):
    add_rect(slide, l, t, w, h, bg_color)
    add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.4,
             font_size=title_size, bold=True, color=title_color)
    tf_top = t + 0.5
    for b in bullets:
        add_text(slide, f"• {b}", l+0.2, tf_top, w-0.4, 0.35,
                 font_size=bullet_size, color=WHITE)
        tf_top += 0.32


def add_flow_arrow(slide, x, y, horizontal=True):
    """Add a simple arrow shape."""
    if horizontal:
        connector = slide.shapes.add_connector(
            1, Inches(x), Inches(y), Inches(x+0.5), Inches(y)
        )
    else:
        connector = slide.shapes.add_connector(
            1, Inches(x), Inches(y), Inches(x), Inches(y+0.3)
        )
    connector.line.color.rgb = YELLOW
    connector.line.width = Pt(2)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)

# accent bar left
add_rect(slide, 0, 0, 0.5, 7.5, HIGHLIGHT)

# title
add_text(slide,
         "Finviz New High Screener",
         1.0, 2.2, 11, 1.2,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# subtitle
add_text(slide,
         "finviz_screener_new_high.py  —  Code Walkthrough",
         1.0, 3.5, 11, 0.6,
         font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# description
add_text(slide,
         "Automatically screens US equities at 52-week highs by combining\n"
         "sector momentum, industry momentum, technical filters, and chart downloads.",
         1.0, 4.3, 11, 1.2,
         font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

add_text(slide, "March 2026", 1.0, 6.8, 4, 0.4, font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 – What the Script Does
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "What the Script Does", 0.3, 0.1, 12, 0.7,
         font_size=26, bold=True, color=WHITE)

add_text(slide,
         "Goal: find US stocks hitting 52-week highs that belong to sectors and "
         "industries with genuine momentum — then download their weekly & monthly charts "
         "for visual review.",
         0.4, 1.1, 12.5, 0.9, font_size=16, color=LIGHT_GRAY)

steps = [
    ("1", "Sector Filter",    "Keep sectors with positive 1M AND 3M performance.\nPick the 'Up Middle' tier (middle 33%).", HIGHLIGHT),
    ("2", "Industry Filter",  "Keep industries with positive 1M performance.",                                               HIGHLIGHT),
    ("3", "Stock Screener",   "Run Finviz screener with 7 technical filters (volume, new high, SMA…).",                     HIGHLIGHT),
    ("4", "Intersect & Rank", "Cross-filter stocks by qualified sectors & industries.\nSummarise counts per group.",         HIGHLIGHT),
    ("5", "Chart Download",   "Download weekly (w) + monthly (m) candlestick charts\nfrom Finviz for each matched stock.",   HIGHLIGHT),
]

for i, (num, title, desc, col) in enumerate(steps):
    x = 0.3 + i * 2.58
    add_rect(slide, x, 2.1, 2.4, 0.55, col)
    add_text(slide, f"Step {num}", x+0.1, 2.15, 2.2, 0.4,
             font_size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_rect(slide, x, 2.65, 2.4, 2.9, ACCENT)
    add_text(slide, title, x+0.1, 2.7, 2.2, 0.45,
             font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x+0.1, 3.2, 2.2, 2.2,
             font_size=12, color=WHITE, align=PP_ALIGN.LEFT)

    if i < 4:
        add_text(slide, "→", x+2.42, 3.3, 0.3, 0.4,
                 font_size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 – Step 1: Sector Filter
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Step 1 — Sector Filter  (The 'Up Middle' Concept)", 0.3, 0.1, 12.5, 0.7,
         font_size=24, bold=True, color=WHITE)

# Left column: logic
add_bullet_box(slide, "Logic",
    [
        "Fetch all sectors via  obb.equity.compare.groups(group='sector')",
        "Clean % strings → float  (clean_performance)",
        "Keep sectors where  performance_1m > 0  AND  performance_3m > 0",
        "Sort by 3M performance (descending)",
        "Slice middle 33 %  →  indices [n*0.33 : n*0.66]",
        "These are the 'Up Middle' qualified sectors",
    ],
    0.3, 1.0, 6.5, 3.5, title_size=15, bullet_size=12)

# Right column: why Up Middle
add_rect(slide, 7.2, 1.0, 5.8, 3.5, ACCENT)
add_text(slide, "Why 'Up Middle'?", 7.4, 1.1, 5.4, 0.4,
         font_size=15, bold=True, color=YELLOW)

tiers = [
    ("Top 33%",    "Already over-extended / late movers",       RED),
    ("Middle 33%", "Momentum building — early-stage leaders ✓", GREEN),
    ("Bottom 33%", "Weak — avoid despite positive perf.",        LIGHT_GRAY),
]
for j, (label, desc, col) in enumerate(tiers):
    y = 1.6 + j * 0.9
    add_rect(slide, 7.4, y, 5.4, 0.75, col)
    add_text(slide, label, 7.5, y+0.05, 1.5, 0.35, font_size=12, bold=True, color=BG_DARK)
    add_text(slide, desc,  9.1, y+0.05, 3.5, 0.35, font_size=12, color=BG_DARK)

# Code snippet
add_rect(slide, 0.3, 4.7, 12.7, 2.4, RGBColor(0x0D, 0x0D, 0x1A))
add_text(slide,
         "df_sectors = obb.equity.compare.groups(group='sector', metric='performance', provider='finviz').to_df()\n"
         "positive_sectors = df_sectors[(df_sectors['performance_1m'] > 0) & (df_sectors['performance_3m'] > 0)]\n"
         "ranked_sectors   = positive_sectors.sort_values(by='performance_3m', ascending=False)\n"
         "n = len(ranked_sectors)\n"
         "up_middle_sectors = ranked_sectors.iloc[int(n*0.33) : int(n*0.66)]\n"
         "qualified_sectors = up_middle_sectors['name'].tolist()",
         0.5, 4.75, 12.3, 2.3, font_size=11, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 – Step 2 & 3: Industry Filter + Screener
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Step 2 & 3 — Industry Filter + Screener Filters", 0.3, 0.1, 12.5, 0.7,
         font_size=24, bold=True, color=WHITE)

# Industry filter
add_bullet_box(slide, "Step 2 — Industry Filter",
    [
        "Fetch all industries via obb.equity.compare.groups(group='industry')",
        "Keep industries where  performance_1m > 0",
        "No ranking — all positive-1M industries qualify",
    ],
    0.3, 1.0, 5.8, 2.5, title_size=14, bullet_size=12)

# Screener filters table
add_rect(slide, 6.5, 1.0, 6.5, 2.5, ACCENT)
add_text(slide, "Step 3 — Screener Filters (Finviz)", 6.6, 1.05, 6.2, 0.4,
         font_size=14, bold=True, color=GREEN)

filters = [
    ("Average Volume",          "Over 300K"),
    ("52-Week High/Low",        "New High"),
    ("Performance",             "Week Up"),
    ("Performance 2",           "Quarter +10%"),
    ("20-Day SMA",              "SMA20 above SMA50"),
    ("200-Day SMA",             "Price above SMA200"),
    ("50-Day SMA",              "SMA50 below SMA20"),
]
for j, (k, v) in enumerate(filters):
    y = 1.55 + j * 0.27
    add_text(slide, k, 6.6, y, 3.0, 0.27, font_size=10, bold=True, color=YELLOW)
    add_text(slide, v, 9.7, y, 3.1, 0.27, font_size=10, color=WHITE)

# Market cap note
add_rect(slide, 0.3, 3.7, 12.7, 0.65, HIGHLIGHT)
add_text(slide,
         "Market Cap filter:  mktcap = 'mid_over'  →  Mid-cap and above only",
         0.5, 3.77, 12.3, 0.5, font_size=13, bold=True, color=WHITE)

# Code snippet
add_rect(slide, 0.3, 4.5, 12.7, 2.65, RGBColor(0x0D, 0x0D, 0x1A))
add_text(slide,
         "# Industry filter\n"
         "qualified_industries = df_industries[df_industries['performance_1m'] > 0]['name'].tolist()\n\n"
         "# Screener\n"
         "results = obb.equity.screener(\n"
         "    provider='finviz',\n"
         "    mktcap='mid_over',\n"
         "    filters_dict=filters   # dict with 7 filters above\n"
         ").to_df()",
         0.5, 4.55, 12.3, 2.55, font_size=11, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 – Step 4: Intersect & Rank
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Step 4 — Intersect Screener Results & Summarise", 0.3, 0.1, 12.5, 0.7,
         font_size=24, bold=True, color=WHITE)

add_bullet_box(slide, "Intersection Filter",
    [
        "Keep stocks where  sector ∈ qualified_sectors",
        "AND  industry ∈ qualified_industries",
        "Result: stocks that pass ALL three gates",
    ],
    0.3, 1.0, 6.0, 2.5, title_size=14, bullet_size=13)

add_bullet_box(slide, "Summary Table",
    [
        "Group by (sector, industry) → count stocks per group",
        "Sort by count descending",
        "Print table for manual review",
        "All industries retained (no further count-cutoff)",
        "Final list sorted by sector (alphabetical)",
    ],
    6.7, 1.0, 6.3, 2.5, title_size=14, bullet_size=13)

# Venn-diagram-like illustration
add_rect(slide, 1.5, 3.7, 4.0, 1.5, RGBColor(0x0F, 0x3C, 0x78))
add_text(slide, "Up-Middle Sectors", 1.6, 3.85, 3.8, 0.5,
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 4.7, 3.7, 4.0, 1.5, RGBColor(0x21, 0x61, 0x80))
add_text(slide, "Positive-1M Industries", 4.8, 3.85, 3.8, 0.5,
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 8.2, 3.7, 4.5, 1.5, RGBColor(0x14, 0x5A, 0x32))
add_text(slide, "Finviz Screener\n(7 filters + mid+ cap)", 8.3, 3.85, 4.2, 0.7,
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "∩", 5.7, 4.0, 0.5, 0.8, font_size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "∩", 9.2, 4.0, 0.5, 0.8, font_size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 5.5, 6.3, 0.7, GREEN)
add_text(slide, "→  Final Candidate Stocks", 3.7, 5.55, 6.0, 0.55,
         font_size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# Code
add_rect(slide, 0.3, 6.3, 12.7, 0.9, RGBColor(0x0D, 0x0D, 0x1A))
add_text(slide,
         "filtered = results[results['sector'].isin(qualified_sectors) & results['industry'].isin(qualified_industries)]\n"
         "summary  = filtered.groupby(['sector', 'industry']).size().reset_index(name='counts').sort_values('counts', ascending=False)",
         0.5, 6.32, 12.3, 0.85, font_size=10, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 – Step 5: Chart Download
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Step 5 — Download Weekly & Monthly Charts", 0.3, 0.1, 12.5, 0.7,
         font_size=24, bold=True, color=WHITE)

add_bullet_box(slide, "Chart Download Logic",
    [
        "Endpoint:  charts2-node.finviz.com/chart.ashx",
        "Params:  cs=m, ct=candle_stick, tm=m  + tf=w or tf=m",
        "Filename pattern:  {Sector}_{Industry}_{Symbol}_{tf}.png",
        "Saved to:  fig/finviz_new_high/",
        "Skip if file already exists (idempotent)",
        "Random sleep 1.0 – 2.5 s between requests (rate-limit avoidance)",
        "On 403 response → pause 60 s, rebuild session, retry once",
    ],
    0.3, 1.0, 7.0, 4.2, title_size=15, bullet_size=12)

add_bullet_box(slide, "Session & Retry Config",
    [
        "requests.Session  with  HTTPAdapter",
        "Retry: total=5, backoff_factor=2",
        "Retry on: 429, 500, 502, 503, 504",
        "Rotating User-Agent pool (5 agents)",
        "Referer header: finviz.com",
        "Timeout: 20 s per request",
    ],
    7.6, 1.0, 5.4, 4.2, title_size=15, bullet_size=12)

# Code
add_rect(slide, 0.3, 5.4, 12.7, 1.8, RGBColor(0x0D, 0x0D, 0x1A))
add_text(slide,
         "for item in tqdm(download_list):\n"
         "    download_chart(symbol, sector, industry, 'w', session)   # weekly\n"
         "    time.sleep(random.uniform(1.0, 2.5))\n"
         "    download_chart(symbol, sector, industry, 'm', session)   # monthly\n"
         "    time.sleep(random.uniform(1.0, 2.5))",
         0.5, 5.45, 12.3, 1.7, font_size=12, color=GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 – Data Flow Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "End-to-End Data Flow", 0.3, 0.1, 12.5, 0.7,
         font_size=26, bold=True, color=WHITE)

nodes = [
    ("Finviz\nSector Perf API",  0.3,  1.3),
    ("Filter +\nUp Middle",      2.5,  1.3),
    ("Finviz\nIndustry Perf API",0.3,  3.4),
    ("Filter 1M > 0",            2.5,  3.4),
    ("Finviz\nScreener API",     0.3,  5.2),
    ("7 Tech Filters\n+ Mid+ Cap",2.5, 5.2),
    ("Intersection",             5.0,  3.4),
    ("Summary\nTable",           7.1,  3.4),
    ("Chart\nDownloader",        9.2,  3.4),
    ("PNG Charts\nsaved locally",11.2, 3.4),
]
colors_n = [HIGHLIGHT, HIGHLIGHT, HIGHLIGHT, HIGHLIGHT,
            HIGHLIGHT, HIGHLIGHT, GREEN,     GREEN,
            YELLOW,   GREEN]

for (label, lx, ty), col in zip(nodes, colors_n):
    add_rect(slide, lx, ty, 1.9, 1.0, col)
    add_text(slide, label, lx+0.05, ty+0.1, 1.8, 0.8,
             font_size=11, bold=True, color=WHITE if col != GREEN else BG_DARK,
             align=PP_ALIGN.CENTER)

# arrows (approximate)
arrows = [
    (2.2,  1.8),   # sector → filter
    (2.2,  3.9),   # industry → filter
    (2.2,  5.7),   # screener → filter
    (4.4,  1.8),   # up middle → intersection (label)
    (4.4,  3.9),   # ind filter → intersection
    (4.4,  5.7),   # tech filter → intersection
    (6.9,  3.9),   # intersection → summary
    (9.0,  3.9),   # summary → downloader
    (11.1, 3.9),   # downloader → png
]
arrow_chars = ["→","→","→","→","→","→","→","→","→"]
positions   = [(2.22,1.7),(2.22,3.7),(2.22,5.5),(4.42,1.7),(4.42,3.7),(4.42,5.5),(6.92,3.7),(9.02,3.7),(11.1,3.7)]
for (ax, ay), ch in zip(positions, arrow_chars):
    add_text(slide, ch, ax, ay, 0.35, 0.4, font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# vertical merge arrows into Intersection box
add_text(slide, "↓", 5.85, 2.35, 0.4, 0.5, font_size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "↑", 5.85, 4.45, 0.4, 0.5, font_size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 – Key Design Decisions & Limitations
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Key Design Decisions & Limitations", 0.3, 0.1, 12.5, 0.7,
         font_size=24, bold=True, color=WHITE)

add_bullet_box(slide, "Design Decisions",
    [
        "OpenBB (obb) as unified data layer — single SDK for sector, industry, and screener data",
        "'Up Middle' avoids chasing extended leaders — targets momentum building stocks",
        "Idempotent chart download — re-runs skip already-downloaded files",
        "Rotating User-Agent + random delays to avoid bot detection",
        "Output directory cleared at start — ensures fresh run each time",
        "Filenames encode sector + industry + symbol → easy sorting/browsing",
    ],
    0.3, 1.0, 6.3, 5.7, title_size=15, bullet_size=12)

add_bullet_box(slide, "Limitations / Watch-outs",
    [
        "Depends on Finviz rate limits — aggressive runs may get 403s",
        "Middle-33% slice is heuristic — may miss leaders in thin sectors",
        "No persistence of screener results (no CSV/DB save step)",
        "No deduplication if output dir is NOT cleared between runs",
        "OpenBB Finviz provider requires valid session / Elite for some data",
        "Charts cleared + re-downloaded every run (no incremental update)",
    ],
    6.8, 1.0, 6.2, 5.7, title_size=15, bullet_size=12, title_color=RED)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 – Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_bg(slide, BG_DARK)
add_rect(slide, 0, 0, 13.33, 0.9, HIGHLIGHT)
add_text(slide, "Summary", 0.3, 0.1, 12.5, 0.7,
         font_size=30, bold=True, color=WHITE)

add_text(slide,
         "finviz_screener_new_high.py  is a 5-step pipeline that applies a sector-momentum "
         "gate, an industry-momentum gate, and 7 technical filters to identify US equities "
         "hitting 52-week highs inside early-stage (Up Middle) sectors — then downloads "
         "weekly and monthly charts for manual review.",
         0.5, 1.1, 12.3, 1.3, font_size=15, color=LIGHT_GRAY)

summary_points = [
    ("Step 1", "Sector Filter",    "Up Middle 33% by 3M performance",       HIGHLIGHT),
    ("Step 2", "Industry Filter",  "All industries with positive 1M return", HIGHLIGHT),
    ("Step 3", "Screener",         "7 filters: volume, new high, SMAs, perf",HIGHLIGHT),
    ("Step 4", "Intersect",        "Stocks passing ALL 3 gates; rank by sector/industry count", HIGHLIGHT),
    ("Step 5", "Charts",           "Weekly + Monthly PNG charts saved locally", HIGHLIGHT),
]

for i, (step, title, desc, col) in enumerate(summary_points):
    y = 2.6 + i * 0.85
    add_rect(slide, 0.5, y, 1.3, 0.7, col)
    add_text(slide, step, 0.5, y+0.1, 1.3, 0.5, font_size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.0, y, 10.8, 0.7, ACCENT)
    add_text(slide, f"{title}  —  {desc}", 2.15, y+0.12, 10.5, 0.5, font_size=13, color=WHITE)

add_text(slide, "Output:  fig/finviz_new_high/{Sector}_{Industry}_{Symbol}_{tf}.png",
         0.5, 7.1, 12.3, 0.35, font_size=13, bold=True, color=GREEN)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "/Users/chihjuihsu/Documents/python_script/TWSectorScreener/finviz_screener_new_high_explained.pptx"
prs.save(out)
print(f"Saved: {out}")
